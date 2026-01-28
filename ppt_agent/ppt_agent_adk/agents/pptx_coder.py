from google.adk.agents import Agent
from ppt_state import PPTState, state_to_dict
import json


def generate_pptx_code(state: PPTState) -> dict:
    """Generates Python code to create the PowerPoint file.
    
    Args:
        state: Current PPT state with all slides and content
        
    Returns:
        dict: Python code to generate the PPTX file
    """
    # Convert state to dictionary for easier processing
    state_dict = state_to_dict(state)
    
    # Generate Python code using python-pptx
    code = f'''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Add slides
    slides_data = {json.dumps(state_dict["slides"], indent=4)}
    
    for slide_data in slides_data:
        # Add slide
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        # Set content
        content = slide.placeholders[1]
        content.text = slide_data["content"]
    
    # Save presentation
    prs.save("{state.file_name}")
    return "{state.file_name}"

# Execute the function
if __name__ == "__main__":
    filename = create_presentation()
    print(f"Presentation saved as: {{filename}}")
'''
    
    return {
        "status": "success",
        "code": code,
        "filename": state.file_name,
        "message": f"Generated code for {state.file_name}"
    }


def execute_pptx_creation(state: PPTState) -> dict:
    """Executes the PowerPoint creation code.
    
    Args:
        state: Current PPT state
        
    Returns:
        dict: Result of PowerPoint creation
    """
    try:
        # Generate the code
        code_result = generate_pptx_code(state)
        
        if code_result["status"] == "success":
            # Execute the code
            exec(code_result["code"])
            
            return {
                "status": "success",
                "filename": code_result["filename"],
                "message": f"Successfully created {code_result['filename']}"
            }
        else:
            return {
                "status": "error",
                "message": "Failed to generate code"
            }
    except Exception as e:
        return {
            "status": "error",
            "message": f"Error creating PowerPoint: {str(e)}"
        }


# Create the PPTX Coder Agent
pptx_coder_agent = Agent(
    name="pptx_coder",
    model="gemini-1.5-flash-002",
    description="Agent responsible for generating and executing PowerPoint creation code.",
    instruction="""You are a PPTX coder. Your job is to:
1. Take the PPT state with slides and content
2. Generate Python code using python-pptx library
3. Execute the code to create the actual PowerPoint file
4. Return the result of the creation process

Always respond with a JSON object containing:
- status: "success" or "error"
- filename: name of the created file
- message: explanation of what was done""",
    tools=[generate_pptx_code, execute_pptx_creation]
) 